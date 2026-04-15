from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "presentation" / "Urban_Climate_Insights_API_Presentation.pptx"
STUDENT_NAME = "Yuan Rui"
STUDENT_ID = "2022115981"
GITHUB_URL = "https://github.com/gapmaker1/urban-climate-insights-api"

NAVY = RGBColor(15, 23, 42)
TEAL = RGBColor(13, 148, 136)
ORANGE = RGBColor(249, 115, 22)
SLATE = RGBColor(51, 65, 85)
LIGHT = RGBColor(241, 245, 249)


def add_title(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.5), Inches(1.1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    subtitle_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(8.0), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(13)
    p.font.color.rgb = SLATE


def add_bullets(slide, title, bullets, accent=TEAL):
    add_title(slide, title, "")
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.2))
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT
    panel.line.color.rgb = accent

    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(8.0), Inches(4.6))
    tf = box.text_frame
    tf.word_wrap = True

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(21 if idx == 0 else 19)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Architecture Overview", "Layered FastAPI design for coursework demonstration and maintenance.")

    labels = [
        ("Client / Swagger UI", Inches(0.8), Inches(2.0), TEAL),
        ("FastAPI Routers", Inches(3.1), Inches(2.0), ORANGE),
        ("Services + Validation", Inches(5.4), Inches(2.0), TEAL),
        ("SQLite + SQLAlchemy", Inches(2.0), Inches(4.1), NAVY),
        ("Open-Meteo APIs", Inches(5.3), Inches(4.1), NAVY),
    ]

    for text, left, top, color in labels:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(2.0), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.text = text
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = NAVY

    for x1, y1, x2, y2 in [
        (1.8, 2.9, 3.1, 2.9),
        (5.1, 2.9, 5.4, 2.9),
        (4.1, 2.9, 4.1, 4.1),
        (6.3, 2.9, 6.3, 4.1),
    ]:
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = SLATE
        line.line.width = Pt(2)

    notes = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(8.3), Inches(1.0))
    tf = notes.text_frame
    p = tf.paragraphs[0]
    p.text = "Key point: writes are authenticated, analytics stay public, and external data import is isolated behind a service layer."
    p.font.size = Pt(17)
    p.font.color.rgb = SLATE


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(title_slide, "Urban Climate Insights API", "Web Services and Web Data coursework presentation")
    hero = title_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.1), Inches(8.7), Inches(2.1))
    hero.fill.solid()
    hero.fill.fore_color.rgb = LIGHT
    hero.line.color.rgb = TEAL
    hero_tf = hero.text_frame
    hero_tf.text = "FastAPI + SQLAlchemy + JWT + Open-Meteo integration"
    hero_tf.paragraphs[0].font.size = Pt(26)
    hero_tf.paragraphs[0].font.bold = True
    hero_tf.paragraphs[0].font.color.rgb = NAVY
    for text, top in [
        (f"Student Name: {STUDENT_NAME}", 4.7),
        (f"Student ID: {STUDENT_ID}", 5.15),
        (f"GitHub: {GITHUB_URL}", 5.6),
    ]:
        box = title_slide.shapes.add_textbox(Inches(0.85), Inches(top), Inches(7.5), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = SLATE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Problem and Aim",
        [
            "Goal: turn open UK city climate data into a reusable web API with full CRUD support.",
            "Users can manage cities and records, import external observations, and run comparative analytics.",
            "The design targets oral demonstration quality as well as technical report depth.",
        ],
    )

    add_architecture_slide(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Data Model and Endpoints",
        [
            "Core models: User, City, UrbanClimateRecord.",
            "CRUD endpoints: /auth, /cities, /records.",
            "Advanced endpoints: /imports/cities/{id}/historical, /analytics/summary, /analytics/compare, /analytics/anomalies.",
        ],
        accent=ORANGE,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Security, Validation, and Errors",
        [
            "JWT bearer authentication protects all write operations.",
            "Pydantic validation enforces field ranges, date logic, and temperature consistency.",
            "HTTP status codes follow convention: 201 created, 204 deleted, 401 unauthorised, 404 not found, 409 conflict, 502 upstream failure.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Documentation, Testing, and Version Control",
        [
            "Swagger UI and exported OpenAPI schema document the live API.",
            "Pytest covers authentication, CRUD, and analytics behaviour.",
            "Technical report highlights stack choice, testing, limitations, and declared GenAI usage.",
            "Repository includes README, PDF API docs, PDF technical report, Dockerfile, and generated slide deck.",
        ],
        accent=ORANGE,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Live Demo Flow",
        [
            "1. Start API and open /docs.",
            "2. Log in with the seeded demo account.",
            "3. Show city list, create a record, then run summary and comparison analytics.",
            "4. Explain how Open-Meteo import populates 30 days of weather and air-quality data.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullets(
        slide,
        "Reflection and Future Work",
        [
            "Strengths: modular code, external data integration, reproducible setup, and coursework-aligned deliverables.",
            "Limitations: SQLite is best for local demo; production deployment should move to PostgreSQL and background jobs.",
            "Future work: caching, scheduled imports, role-based access control, and visual dashboards.",
        ],
        accent=ORANGE,
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Presentation written to {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
